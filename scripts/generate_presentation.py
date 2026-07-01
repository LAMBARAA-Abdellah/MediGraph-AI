"""Generate MediGraph AI PowerPoint presentation (20-25 slides)."""

from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("Install python-pptx: pip install python-pptx")
    raise

OUTPUT = Path(__file__).parent.parent / "docs" / "reports" / "MediGraph_AI_Presentation.pptx"

SLIDES = [
    ("MediGraph AI", "Intelligent Multi-Agent Clinical Orientation System\nUniversité Mohammed VI Polytechnique — Master BDCC"),
    ("Agenda", "• Introduction\n• Architecture\n• Clinical Workflow\n• LangGraph Multi-Agent System\n• MCP Integration\n• FastAPI Backend\n• React Frontend\n• Testing\n• Conclusion"),
    ("Introduction", "Educational multi-agent clinical orientation platform\n\n⚠ Does NOT provide medical diagnosis\n⚠ Does NOT replace healthcare professionals\n\nDisclaimer: This system does not replace a professional medical consultation."),
    ("Project Objectives", "• Simulate clinical workflow with AI agents\n• Supervisor, Diagnostic, Physician Review, Report agents\n• Human-in-the-Loop physician review\n• MCP tool integration\n• Multi-format report generation\n• Production-ready architecture"),
    ("Technologies", "• Python 3.12\n• LangGraph & LangChain\n• FastAPI\n• MCP (Model Context Protocol)\n• React Frontend\n• SQLite & SQLAlchemy\n• Docker & Docker Compose\n• LangGraph Studio"),
    ("System Architecture", "Four-layer clean architecture:\n\n1. Presentation — React + FastAPI\n2. Application — Services + LangGraph\n3. Domain — State, Nodes, Tools\n4. Infrastructure — DB, MCP, Redis"),
    ("Architecture Diagram", "React Frontend → FastAPI → LangGraph\n                                    ↓\n                              MCP Server + SQLite\n\nModular, scalable, SOLID principles"),
    ("Clinical Workflow", "1. Patient Registration\n2. 5 Sequential Questions\n3. Clinical Summary Generation\n4. Intermediate Recommendation\n5. Physician Review (HITL)\n6. Final Report (JSON/MD/HTML/PDF)"),
    ("Workflow Diagram", "Supervisor → Diagnostic Agent → Physician Review → Report Agent\n\nConditional routing based on consultation_status\nInterrupt before physician_review for HITL"),
    ("LangGraph State", "MedicalState TypedDict:\n• messages (MessagesState)\n• patient_information\n• question_count / patient_answers\n• clinical_summary\n• intermediate_recommendation\n• physician_review / physician_treatment\n• final_report / consultation_status"),
    ("Supervisor Agent", "• Workflow orchestration\n• Decision making & node routing\n• Error handling\n• Conditional graph transitions\n• Entry point of the graph"),
    ("Diagnostic Agent", "• Asks exactly 5 sequential questions\n• Questions via ask_patient tool\n• Collects answers in shared state\n• Generates Clinical Summary\n• Generates Intermediate Recommendation\n• Never provides diagnosis"),
    ("Physician Review Agent", "Human-in-the-Loop:\n• Workflow PAUSES at interrupt\n• Physician reviews summary\n• Can modify recommendation\n• Add treatment notes\n• Approve or reject\n• Resume via API"),
    ("Report Agent", "Generates structured reports:\n• JSON output\n• Markdown report\n• HTML report (Jinja2)\n• PDF report (ReportLab)\n\nMandatory disclaimer on every report"),
    ("MCP Integration", "MCP Server tools:\n• Patient tools\n• Recommendation tools\n• Medical knowledge lookup\n• Drug database demo\n• Calendar/appointment\n\nClient with local fallback"),
    ("FastAPI Backend", "Endpoints:\n• POST /sessions/start\n• POST /consultation/start\n• POST /consultation/answer\n• POST /consultation/physician-review\n• GET /consultation/{id}/report\n• GET /health, /metrics"),
    ("React Frontend", "Pages:\n• Dashboard\n• Patient Registration\n• Consultation & Questionnaire\n• Physician Review\n• Reports & History\n• Settings"),
    ("Frontend Screens", "Patient Screen: Name, Age, Gender, History, Complaint\nQuestion Screen: Progress bar, Q&A\nPhysician Screen: Summary, Approve/Reject\nReport Screen: Full report + PDF export"),
    ("Database Design", "SQLite tables:\n• Patient\n• Consultation\n• Question / Answer\n• PhysicianReview\n• Report\n\nLangGraph checkpoint persistence"),
    ("LangGraph Studio", "Configuration: langgraph.json\n• Graph visualization\n• State inspection\n• Checkpoint debugging\n• Interrupt point testing\n\nRun: langgraph dev"),
    ("Testing", "Case 1: Simple Respiratory Syndrome\nCase 2: Red Flags (high urgency)\nCase 3: Benign Case (routine)\n\nVerify: 5 questions, recommendations, HITL, reports, disclaimer"),
    ("Docker Deployment", "docker-compose up --build\n\nServices:\n• Backend (port 8000)\n• MCP Server (port 8100)\n• Frontend (port 5173)\n• Redis (optional)"),
    ("Results", "✅ Complete multi-agent workflow\n✅ 5 sequential questions\n✅ Clinical Summary & Recommendation\n✅ Human-in-the-Loop review\n✅ Multi-format reports\n✅ MCP integration\n✅ Full stack deployment"),
    ("Conclusion", "MediGraph AI demonstrates production-ready multi-agent clinical orientation using LangGraph, LangChain, FastAPI, and MCP.\n\nEducational project for academic purposes only."),
    ("Future Improvements", "• EHR integration\n• Multi-language support\n• Advanced LLM summaries\n• Kubernetes deployment\n• WebSocket real-time updates\n• OAuth2 security\n• FHIR compliance"),
    ("Thank You", "MediGraph AI\n\nQuestions?\n\nDisclaimer: This system does not replace a professional medical consultation."),
]


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for title, content in SLIDES:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()

        for i, line in enumerate(content.split("\n")):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0x2D, 0x37, 0x48)

        # Style title
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x2C, 0x52, 0x82)
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Presentation saved to: {OUTPUT}")
    print(f"Total slides: {len(SLIDES)}")


if __name__ == "__main__":
    create_presentation()
